from __future__ import annotations

import csv
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent
FINAL_DIR = BASE / "最终答辩材料"
OUT = FINAL_DIR / "相机知识图谱期末答辩PPT_最终版.pptx"
ROOT_OUT = BASE / OUT.name

NEO4J_URL = "http://localhost:7474/browser/"
ADVISOR_URL = "http://127.0.0.1:5000"

BG = RGBColor(247, 250, 252)
NAVY = RGBColor(22, 45, 64)
BLUE = RGBColor(35, 116, 171)
TEAL = RGBColor(29, 158, 147)
GREEN = RGBColor(66, 153, 111)
AMBER = RGBColor(224, 153, 54)
RED = RGBColor(214, 90, 90)
PURPLE = RGBColor(113, 91, 180)
TEXT = RGBColor(43, 52, 62)
MUTED = RGBColor(98, 112, 126)
LINE = RGBColor(214, 225, 235)
CODE_BG = RGBColor(25, 35, 46)
WHITE = RGBColor(255, 255, 255)
FONT_CN = "Microsoft YaHei"
FONT_MONO = "Consolas"


def count_csv() -> tuple[int, int, dict[str, int], dict[str, int]]:
    node_counts: dict[str, int] = {}
    edge_counts: dict[str, int] = {}
    node_total = edge_total = 0
    with (FINAL_DIR / "normalized_nodes.csv").open("r", encoding="utf-8-sig", newline="") as f:
        for row in csv.DictReader(f):
            node_total += 1
            label = row.get("Label", "")
            node_counts[label] = node_counts.get(label, 0) + 1
    with (FINAL_DIR / "normalized_edges.csv").open("r", encoding="utf-8-sig", newline="") as f:
        for row in csv.DictReader(f):
            edge_total += 1
            rel_type = row.get("type", "")
            edge_counts[rel_type] = edge_counts.get(rel_type, 0) + 1
    return node_total, edge_total, node_counts, edge_counts


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
NODE_TOTAL, EDGE_TOTAL, NODE_COUNTS, EDGE_COUNTS = count_csv()


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def text(slide, x, y, w, h, content, size=16, color=TEXT, bold=False, align=PP_ALIGN.LEFT, font=FONT_CN):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = content
    p.alignment = align
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return shape


def title(slide, title_text, subtitle=None, section=None):
    x = 0.55
    if section:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.32), Inches(1.25), Inches(0.34))
        pill.fill.solid()
        pill.fill.fore_color.rgb = TEAL
        pill.line.fill.background()
        pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = pill.text_frame.paragraphs[0]
        p.text = section
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = FONT_CN
        p.runs[0].font.size = Pt(10)
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.bold = True
        x = 1.95
    text(slide, x, 0.25, 9.9, 0.5, title_text, 24, NAVY, True)
    if subtitle:
        text(slide, x, 0.72, 10.8, 0.28, subtitle, 11, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.08), Inches(12.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def footer(slide, idx):
    text(slide, 10.7, 7.12, 2.0, 0.22, f"{idx:02d} / 13", 9, MUTED, align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, h, items, size=14):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(6)
        for run in p.runs:
            run.font.name = FONT_CN
            run.font.size = Pt(size)
            run.font.color.rgb = TEXT
    return shape


def card(slide, x, y, w, h, heading, items=None, color=BLUE, icon=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE
    if icon:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.42), Inches(0.42))
        c.fill.solid()
        c.fill.fore_color.rgb = color
        c.line.fill.background()
        text(slide, x + 0.18, y + 0.29, 0.42, 0.18, icon, 9, WHITE, True, PP_ALIGN.CENTER)
        tx = x + 0.72
    else:
        tx = x + 0.22
    text(slide, tx, y + 0.17, w - (tx - x) - 0.18, 0.34, heading, 15, NAVY, True)
    if items:
        bullets(slide, x + 0.22, y + 0.66, w - 0.44, h - 0.78, items, 12)
    return box


def metric(slide, x, y, label, value, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.25), Inches(0.95))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE
    text(slide, x + 0.15, y + 0.12, 1.9, 0.28, str(value), 22, color, True, PP_ALIGN.CENTER)
    text(slide, x + 0.15, y + 0.55, 1.9, 0.22, label, 10, MUTED, align=PP_ALIGN.CENTER)


def code(slide, x, y, w, h, content, size=12):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = RGBColor(68, 82, 96)
    shape = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.16), Inches(w - 0.36), Inches(h - 0.32))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    for run in p.runs:
        run.font.name = FONT_MONO
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(218, 232, 240)
    return box


def button(slide, x, y, w, h, label, url, color=TEAL):
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    btn.fill.solid()
    btn.fill.fore_color.rgb = color
    btn.line.fill.background()
    btn.click_action.hyperlink.address = url
    btn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = btn.text_frame.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.name = FONT_CN
    p.runs[0].font.size = Pt(14)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE
    return btn


def flow(slide, x, y, labels, colors):
    box_w, box_h, gap = 1.85, 0.72, 0.25
    for i, label in enumerate(labels):
        bx = x + i * (box_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(y), Inches(box_w), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = box.text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = FONT_CN
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        if i < len(labels) - 1:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(bx + box_w + 0.02),
                Inches(y + box_h / 2),
                Inches(bx + box_w + gap - 0.03),
                Inches(y + box_h / 2),
            )
            conn.line.color.rgb = LINE
            conn.line.width = Pt(2)


def mini_graph(slide, x, y):
    nodes = [
        (x + 0.5, y + 0.7, "品牌", BLUE),
        (x + 2.1, y + 0.45, "型号", TEAL),
        (x + 3.7, y + 0.75, "镜头", GREEN),
        (x + 1.2, y + 2.1, "用途", AMBER),
        (x + 2.8, y + 2.15, "价格", RED),
        (x + 4.4, y + 2.0, "评测", PURPLE),
    ]
    for a, b in [(0, 1), (1, 2), (1, 3), (1, 4), (1, 5), (2, 3), (4, 5)]:
        xa, ya, _, _ = nodes[a]
        xb, yb, _, _ = nodes[b]
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(xa + 0.35), Inches(ya + 0.18), Inches(xb + 0.35), Inches(yb + 0.18))
        conn.line.color.rgb = RGBColor(180, 195, 210)
        conn.line.width = Pt(1.6)
    for nx, ny, label, color in nodes:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(nx), Inches(ny), Inches(0.75), Inches(0.75))
        c.fill.solid()
        c.fill.fore_color.rgb = color
        c.line.color.rgb = WHITE
        c.line.width = Pt(1.5)
        text(slide, nx, ny + 0.25, 0.75, 0.18, label, 9, WHITE, True, PP_ALIGN.CENTER)


def slide(title_text, subtitle=None, section=None):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    idx = len(prs.slides)
    title(s, title_text, subtitle, section)
    footer(s, idx)
    return s


# 1
s = prs.slides.add_slide(BLANK)
bg(s, RGBColor(239, 247, 250))
text(s, 0.8, 1.0, 9.8, 0.75, "面向消费决策的相机领域知识图谱设计与导购系统实现", 28, NAVY, True)
text(s, 0.85, 1.86, 7.6, 0.34, "知识图谱期末大作业 · Neo4j · Flask · 自然语言导购", 15, MUTED)
flow(s, 0.9, 4.9, ["文献综述", "图谱构建", "Neo4j 查询", "导购系统"], [BLUE, TEAL, GREEN, AMBER])
mini_graph(s, 7.1, 2.0)
text(s, 0.9, 6.85, 3.8, 0.25, "汇报时长：10 分钟", 11, MUTED)
footer(s, 1)

# 2
s = slide("研究背景与问题", "相机选购是一个典型的多实体、多关系、多约束决策场景", "背景")
card(s, 0.75, 1.45, 3.7, 1.25, "信息分散", ["价格、参数、评测分布在电商、官网和媒体平台", "用户需要跨来源比较"], BLUE, "1")
card(s, 4.8, 1.45, 3.7, 1.25, "关系复杂", ["相机与品牌、画幅、镜头、用途、参数相互关联", "表格不容易表达兼容路径"], TEAL, "2")
card(s, 8.85, 1.45, 3.7, 1.25, "推荐需要解释", ["用户关心为什么适合自己", "需要结合预算、用途和已有镜头"], AMBER, "3")
text(s, 0.85, 3.35, 2.0, 0.3, "核心问题", 17, NAVY, True)
bullets(s, 1.05, 3.82, 5.2, 1.9, ["如何把碎片化的相机信息组织成结构化知识？", "如何支持按预算、用途、卡口、镜头生态进行查询？", "如何让推荐结果具有可解释性？"], 16)
mini_graph(s, 6.8, 3.25)

# 3
s = slide("文献综述介绍", "30 篇以上文献按主题归纳，而不是逐篇堆砌", "综述")
card(s, 0.75, 1.45, 3.8, 3.7, "知识图谱构建", ["实体、关系、属性、本体建模", "图数据库存储与图谱质量控制", "支撑本项目的 Schema 设计"], BLUE, "KG")
card(s, 4.85, 1.45, 3.8, 3.7, "推荐系统", ["用户画像、场景匹配、约束过滤", "可解释推荐与冷启动问题", "支撑预算、用途、偏好建模"], TEAL, "RS")
card(s, 8.95, 1.45, 3.8, 3.7, "智能问答", ["自然语言理解与意图识别", "问题到结构化查询的映射", "支撑自然语言转 Cypher"], PURPLE, "QA")
text(s, 0.95, 5.75, 11.3, 0.46, "综述结论：知识图谱适合相机领域的关系表达，可进一步服务于可解释导购与问答应用。", 17, NAVY, True, PP_ALIGN.CENTER)

# 4
s = slide("文献综述对项目设计的启发", "把文献方法落到相机选购场景中的实体、关系和应用能力", "综述")
flow(s, 0.95, 1.55, ["文献方法", "领域建模", "图数据库", "导购应用"], [BLUE, TEAL, GREEN, AMBER])
card(s, 0.9, 2.7, 3.75, 2.1, "图谱构建启发", ["相机、镜头、品牌作为实体", "生产、分类、兼容、适用作为关系"], BLUE)
card(s, 4.8, 2.7, 3.75, 2.1, "推荐系统启发", ["引入预算、用途、经验、便携需求", "把推荐理由显式建成节点/属性"], TEAL)
card(s, 8.7, 2.7, 3.75, 2.1, "问答系统启发", ["自然语言问题解析为约束条件", "生成 Cypher 查询 Neo4j 图谱"], PURPLE)
text(s, 1.0, 5.55, 11.1, 0.5, "因此，本项目不是单纯罗列相机参数，而是把“知识组织”和“购买建议”结合起来。", 18, NAVY, True, PP_ALIGN.CENTER)

# 5
s = slide("系统总体架构", "数据层、图谱层、应用层构成完整闭环", "架构")
flow(s, 0.8, 1.45, ["数据来源", "预处理", "Schema", "Neo4j 导入", "Cypher 查询", "导购系统"], [BLUE, TEAL, GREEN, AMBER, PURPLE, RED])
card(s, 0.85, 2.65, 3.45, 2.2, "数据层", ["品牌/型号/镜头/参数", "价格、评测、用户画像", "归一化 CSV 节点表与边表"], BLUE)
card(s, 4.95, 2.65, 3.45, 2.2, "图谱层", ["Neo4j 存储实体与关系", "支持路径查询与关系解释", "展示整体知识网络"], TEAL)
card(s, 9.05, 2.65, 3.45, 2.2, "应用层", ["自然语言解析", "自动生成 Cypher", "输出推荐结果与理由"], AMBER)
button(s, 3.1, 5.55, 2.25, 0.55, "Neo4j 图谱库", NEO4J_URL, BLUE)
button(s, 7.75, 5.55, 2.55, 0.55, "Flask 导购系统", ADVISOR_URL, TEAL)

# 6
s = slide("数据来源", "围绕市场主流相机、镜头和购买决策信息进行扩展", "数据")
card(s, 0.75, 1.35, 3.05, 1.55, "品牌与型号", ["10 个品牌", "每个品牌至少 15 个热门型号"], BLUE)
card(s, 3.95, 1.35, 3.05, 1.55, "价格信息", ["京东、天猫、品牌官网", "首发价、当前价、二手价"], TEAL)
card(s, 7.15, 1.35, 2.75, 1.55, "参数信息", ["像素、防抖、对焦点", "续航、卡槽、屏幕、码率"], GREEN)
card(s, 10.05, 1.35, 2.55, 1.55, "评测信息", ["媒体评分、用户评分", "优点、缺点、发布日期"], PURPLE)
card(s, 1.15, 3.55, 4.95, 1.7, "镜头生态", ["卡口、焦段、光圈、重量、防抖、价格", "补充市场上热销 30 款镜头"], AMBER)
card(s, 7.0, 3.55, 4.95, 1.7, "用户画像与用途", ["预算、已有镜头、拍摄经验、主要题材、便携要求", "风光、人像、视频、直播、街拍、体育、生态、商业棚拍"], RED)

# 7
s = slide("数据预处理过程", "把分散商品信息转换成统一的节点和关系", "数据")
flow(s, 0.9, 1.45, ["采集整理", "字段清洗", "实体归一", "关系补全", "CSV 输出"], [BLUE, TEAL, GREEN, AMBER, PURPLE])
card(s, 0.85, 2.65, 3.65, 2.4, "清洗与统一", ["统一品牌名称、型号名称、卡口名称", "补齐价格、重量、用途、参数字段", "处理重复和缺失信息"], BLUE)
card(s, 4.85, 2.65, 3.65, 2.4, "节点归一化", ["输出 normalized_nodes.csv", "每一行对应一个实体节点", "用 Label 区分节点类型"], TEAL)
card(s, 8.85, 2.65, 3.65, 2.4, "关系归一化", ["输出 normalized_edges.csv", "每一行对应一条关系边", "用 type 区分关系语义"], AMBER)
text(s, 1.2, 5.65, 10.9, 0.38, "预处理目标：让 Neo4j 导入后形成稳定、可查询、可展示的领域知识网络。", 17, NAVY, True, PP_ALIGN.CENTER)

# 8
s = slide("Schema 设计", "围绕“相机是什么、适合谁、配什么、为什么推荐”建模", "Schema")
card(s, 0.75, 1.35, 5.55, 4.85, "主要节点类型", ["Brand：品牌", "CameraModel：相机型号", "CameraCategory：相机类别", "Lens：镜头", "UseCase：用途场景", "PriceRecord：价格记录", "ReviewRecord：评测记录", "UserProfile：用户画像"], BLUE)
card(s, 7.05, 1.35, 5.55, 4.85, "主要关系类型", ["PRODUCED_BY：生产关系", "BELONGS_TO_CATEGORY：类别归属", "COMPATIBLE_WITH：镜头兼容", "SUITABLE_FOR：相机适用场景", "LENS_SUITABLE_FOR：镜头适用场景", "PRICE_OF / REVIEW_OF：价格与评测关联", "HAS_REASON：推荐理由"], TEAL)

# 9
s = slide("数据库图谱整体展示", "点击按钮进入 Neo4j Browser，现场展示图谱结构", "Neo4j")
metric(s, 0.8, 1.35, "节点总数", NODE_TOTAL, BLUE)
metric(s, 3.25, 1.35, "关系总数", EDGE_TOTAL, TEAL)
metric(s, 5.7, 1.35, "相机型号", NODE_COUNTS.get("CameraModel", 0), GREEN)
metric(s, 8.15, 1.35, "镜头", NODE_COUNTS.get("Lens", 0), AMBER)
metric(s, 10.6, 1.35, "品牌", NODE_COUNTS.get("Brand", 0), PURPLE)
mini_graph(s, 0.95, 3.0)
code(s, 6.55, 2.75, 5.8, 1.4, "MATCH p=()-[]->()\nRETURN p\nLIMIT 120;", 16)
button(s, 7.45, 4.65, 4.0, 0.65, "打开 Neo4j 知识图谱", NEO4J_URL, BLUE)
text(s, 6.65, 5.58, 5.5, 0.42, "答辩现场：点击按钮 → 登录 Neo4j → 运行上方 Cypher → 展示整体图谱。", 12, MUTED, align=PP_ALIGN.CENTER)

# 10
s = slide("Neo4j 图谱可视化展示", "通过节点与关系路径观察相机领域知识结构", "Neo4j")
card(s, 0.75, 1.35, 3.6, 1.55, "可视化观察点", ["品牌连接多个型号", "机身通过兼容关系连接镜头", "用途节点连接推荐场景"], BLUE)
card(s, 0.75, 3.15, 3.6, 1.55, "讲解路径示例", ["索尼 → Alpha 7 IV", "Alpha 7 IV → 全画幅微单", "Alpha 7 IV → FE 镜头"], TEAL)
mini_graph(s, 5.0, 1.5)
code(s, 5.15, 5.05, 6.5, 0.9, "MATCH (m:CameraModel)-[:COMPATIBLE_WITH]->(l:Lens)\nRETURN m,l LIMIT 50;", 12)

# 11
s = slide("功能查询展示", "从图谱关系出发，支撑相机导购中的多类问题", "查询")
card(s, 0.65, 1.28, 3.0, 1.15, "1. 最新机型", ["按 releaseYear 倒序查询", "回答“最新/新款相机”"], BLUE)
card(s, 3.88, 1.28, 3.0, 1.15, "2. 新手推荐", ["用途=入门学习", "结合预算和热销偏好"], TEAL)
card(s, 7.11, 1.28, 3.0, 1.15, "3. 视频创作", ["用途=视频创作", "筛选品牌、预算和推荐分"], GREEN)
card(s, 10.34, 1.28, 2.35, 1.15, "4. 镜头推荐", ["按卡口/用途/价格", "推荐人像或风光镜头"], AMBER)
card(s, 0.65, 2.68, 3.0, 1.15, "5. 兼容机身", ["已有某品牌镜头", "查询可兼容机身"], PURPLE)
card(s, 3.88, 2.68, 3.0, 1.15, "6. 型号解释", ["查询品牌、类别、理由", "展示兼容镜头路径"], RED)
card(s, 7.11, 2.68, 5.58, 1.15, "答辩展示逻辑", ["先展示 Schema，再用功能查询证明图谱可服务购买决策", "查询结果不只是表格，还能解释“为什么推荐”"], BLUE)
representative = """// 示例：已有索尼镜头时推荐可兼容机身
MATCH (m:CameraModel)-[:COMPATIBLE_WITH]->(l:Lens)
      -[:PRODUCED_BY]->(:Brand {name:'索尼'})
RETURN m.name AS 推荐机身, m.price AS 价格,
       m.releaseYear AS 发布年份,
       count(DISTINCT l) AS 可兼容镜头数,
       collect(DISTINCT l.name)[0..5] AS 可用镜头
ORDER BY 可兼容镜头数 DESC, 发布年份 DESC, 价格 ASC
LIMIT 8;"""
code(s, 0.75, 4.18, 7.05, 1.95, representative, 10)
card(s, 8.15, 4.18, 4.35, 1.95, "现场可演示问题", ["最新的相机推荐", "新手预算5000以内，推荐热销相机", "索尼E卡口人像镜头推荐，预算5000以内", "已有索尼镜头，应该买什么机身？"], TEAL)
text(s, 0.95, 6.38, 11.5, 0.32, "展示重点：预算筛选、用途匹配、镜头生态、兼容关系、型号解释、最新排序。", 15, NAVY, True, PP_ALIGN.CENTER)

# 12
s = slide("对话导购系统展示", "自然语言问题自动映射为 Cypher 查询，并返回推荐结果", "应用")
button(s, 0.9, 1.35, 3.45, 0.65, "打开相机导购对话系统", ADVISOR_URL, TEAL)
card(s, 0.9, 2.25, 3.75, 2.8, "现场示例问题", ["新手预算5000以内，推荐热销相机", "索尼E卡口人像镜头推荐，预算5000以内", "已有索尼镜头，应该买什么机身？"], TEAL)
flow(s, 5.25, 1.45, ["用户问题", "语义解析", "Cypher", "Neo4j", "推荐结果"], [BLUE, TEAL, GREEN, AMBER, RED])
card(s, 5.25, 2.65, 3.0, 2.25, "识别字段", ["预算、用途、品牌、类别", "卡口、已有镜头、具体型号", "热销/性价比偏好"], BLUE)
card(s, 8.75, 2.65, 3.2, 2.25, "输出内容", ["推荐型号/镜头", "价格、评分、参数", "推荐理由与生成的 Cypher"], AMBER)
text(s, 5.35, 5.55, 6.45, 0.38, "无法理解的问题会提示：“暂时没能理解您的意思，请重新输入”。", 13, MUTED, align=PP_ALIGN.CENTER)

# 13
s = slide("总结与展望：知识图谱 + 大模型导购", "从结构化知识组织，进一步走向更自然、更智能的消费决策助手", "总结")
card(s, 0.75, 1.25, 3.75, 2.15, "已完成成果", ["完成 30 篇以上文献综述", f"构建 {NODE_TOTAL} 个节点、{EDGE_TOTAL} 条关系", "实现 Neo4j 图谱展示、Cypher 查询和 Flask 导购系统"], BLUE)
card(s, 4.8, 1.25, 3.75, 2.15, "当前系统价值", ["把相机参数、镜头生态、价格和评测结构化", "支持预算、用途、卡口、已有镜头等导购查询", "推荐结果可追溯到图谱关系和理由"], TEAL)
card(s, 8.85, 1.25, 3.75, 2.15, "现有不足", ["规则式语义解析覆盖有限", "价格和销量更新仍依赖人工或定期采集", "复杂偏好和多轮追问能力还不够强"], AMBER)
card(s, 0.75, 3.85, 3.75, 1.85, "大模型改进 1：语义理解", ["用大模型识别隐含需求，例如“拍娃”“旅游轻便”“预算别太高”", "将自然语言稳定转换为 Cypher 或结构化查询参数"], PURPLE)
card(s, 4.8, 3.85, 3.75, 1.85, "大模型改进 2：多轮导购", ["根据用户追问持续更新画像", "支持“再便宜点”“更轻一点”“换佳能”这类上下文指令"], RED)
card(s, 8.85, 3.85, 3.75, 1.85, "大模型改进 3：解释生成", ["结合图谱路径生成更自然的推荐理由", "把参数、评测和价格转化为用户能理解的购买建议"], GREEN)
text(s, 0.95, 6.35, 11.5, 0.34, "未来方向：以知识图谱保证事实可靠，以大模型提升交互体验和推荐解释能力。", 16, NAVY, True, PP_ALIGN.CENTER)


try:
    prs.save(OUT)
    prs.save(ROOT_OUT)
    print(OUT)
    print(ROOT_OUT)
except PermissionError:
    stamp = datetime.now().strftime("%H%M%S")
    alt_out = FINAL_DIR / f"相机知识图谱期末答辩PPT_大模型展望增强版_{stamp}.pptx"
    alt_root_out = BASE / alt_out.name
    prs.save(alt_out)
    prs.save(alt_root_out)
    print(alt_out)
    print(alt_root_out)
